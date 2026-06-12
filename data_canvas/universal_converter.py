"""Universal Data Converter - Convert any raw data format to PowerPoint with performance optimization and structured summarization"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import io
import json
from pathlib import Path
from datetime import datetime
from .utils import parse_raw_data, get_numeric_columns, get_categorical_columns
from .analyzers import DataAnalyzer
from .generators import ChartGenerator
#from .performance import PerformanceOptimizer, ChunkProcessor, CacheManager
#from .summarizer import DataSummarizer, StructuredReportGenerator, PerformanceSummary
import pandas as pd
import time


class UniversalDataToPPT:
    """Convert any raw data format (text, CSV, JSON, Excel, etc.) to PowerPoint with performance optimization"""
    
    def __init__(self, title="Data Analysis Report", auto_format=True, enable_optimization=True):
        """
        Initialize the universal converter.
        
        Args:
            title (str): Presentation title
            auto_format (bool): Automatically detect and format data
            enable_optimization (bool): Enable performance optimization for large datasets
        """
        self.title = title
        self.prs = None
        self.chart_generator = ChartGenerator()
        self.auto_format = auto_format
        self.data_format = None
        self.enable_optimization = enable_optimization
       # self.performance_optimizer = PerformanceOptimizer()
        #self.cache_manager = CacheManager()
       # self.chunk_processor = ChunkProcessor()
       # self.report_generator = StructuredReportGenerator()
        self.conversion_start_time = None
        self.structured_report = None
    
    def convert_text(self, raw_text, output_file, delimiter=None, generate_report=False, report_file=None):
        """
        Convert raw text data to PowerPoint.
        Supports: CSV, TSV, pipe-delimited, semicolon-delimited, space-separated
        
        Args:
            raw_text (str): Raw text data
            output_file (str): Path to output PPT file
            delimiter (str): Data delimiter (auto-detect if None)
            generate_report (bool): Generate structured JSON report
            report_file (str): Path to save JSON report
        """
        try:
            self.conversion_start_time = time.time()
            df = parse_raw_data(raw_text, delimiter)
            return self._convert_dataframe(df, output_file, generate_report, report_file)
        except Exception as e:
            print(f"✗ Error converting text data: {e}")
            return False
    
    def convert_file(self, input_file, output_file, generate_report=False, report_file=None):
        """
        Convert data file to PowerPoint.
        Supports: CSV, Excel, JSON, TXT, TSV, DAT
        
        Args:
            input_file (str): Path to input file
            output_file (str): Path to output PPT file
            generate_report (bool): Generate structured JSON report
            report_file (str): Path to save JSON report
        """
        try:
            self.conversion_start_time = time.time()
            file_path = Path(input_file)
            
            if not file_path.exists():
                raise FileNotFoundError(f"File not found: {input_file}")
            
            extension = file_path.suffix.lower()
            
            if extension == '.csv':
                df = pd.read_csv(input_file)
                self.data_format = 'CSV File'
            elif extension in ['.xls', '.xlsx']:
                df = pd.read_excel(input_file)
                self.data_format = 'Excel File'
            elif extension == '.json':
                df = pd.read_json(input_file)
                self.data_format = 'JSON File'
            elif extension in ['.txt', '.tsv', '.dat']:
                with open(input_file, 'r', encoding='utf-8') as f:
                    raw_text = f.read()
                df = parse_raw_data(raw_text)
                self.data_format = 'Text File'
            else:
                with open(input_file, 'r', encoding='utf-8') as f:
                    raw_text = f.read()
                df = parse_raw_data(raw_text)
                self.data_format = 'Text File'
            
            return self._convert_dataframe(df, output_file, generate_report, report_file)
            
        except Exception as e:
            print(f"✗ Error converting file: {e}")
            return False
    
    def convert_dict(self, data_dict, output_file, generate_report=False, report_file=None):
        """
        Convert Python dictionary to PowerPoint.
        
        Args:
            data_dict (dict): Dictionary with lists or similar structure
            output_file (str): Path to output PPT file
            generate_report (bool): Generate structured JSON report
            report_file (str): Path to save JSON report
        """
        try:
            self.conversion_start_time = time.time()
            df = pd.DataFrame(data_dict)
            self.data_format = 'Python Dictionary'
            return self._convert_dataframe(df, output_file, generate_report, report_file)
        except Exception as e:
            print(f"✗ Error converting dictionary: {e}")
            return False
    
    def convert_json_string(self, json_string, output_file, generate_report=False, report_file=None):
        """
        Convert JSON string to PowerPoint.
        
        Args:
            json_string (str): JSON formatted string
            output_file (str): Path to output PPT file
            generate_report (bool): Generate structured JSON report
            report_file (str): Path to save JSON report
        """
        try:
            self.conversion_start_time = time.time()
            data = json.loads(json_string)
            df = pd.DataFrame(data)
            self.data_format = 'JSON String'
            return self._convert_dataframe(df, output_file, generate_report, report_file)
        except Exception as e:
            print(f"✗ Error converting JSON string: {e}")
            return False
    
    def convert_dataframe(self, df, output_file, generate_report=False, report_file=None):
        """
        Convert pandas DataFrame to PowerPoint.
        
        Args:
            df (pd.DataFrame): Input dataframe
            output_file (str): Path to output PPT file
            generate_report (bool): Generate structured JSON report
            report_file (str): Path to save JSON report
        """
        try:
            self.conversion_start_time = time.time()
            self.data_format = 'Pandas DataFrame'
            return self._convert_dataframe(df, output_file, generate_report, report_file)
        except Exception as e:
            print(f"✗ Error converting DataFrame: {e}")
            return False
    
    def _convert_dataframe(self, df, output_file, generate_report=False, report_file=None):
        """Internal method to convert dataframe"""
        # Generate structured report if requested
        if generate_report:
            self.structured_report = self.report_generator.generate_report(df, self.title)
            if report_file:
                self.report_generator.export_report(report_file)
                print(f"✓ Report exported to: {report_file}")
        
        # Optimize dataframe if enabled
       # if self.enable_optimization:
          #  df = self.performance_optimizer.optimize_dataframe(df)
        
        # Create presentation
        self._create_presentation(df)
        self.prs.save(output_file)
        
        # Calculate and log performance metrics
       # elapsed_time = time.time() - self.conversion_start_time
      #  perf_summary = PerformanceSummary.summarize_conversion(
          #  input_size_mb=df.memory_usage(deep=True).sum() / 1024 / 1024,
         #   num_rows=len(df),
         #   num_columns=len(df.columns),
          #  processing_time_sec=elapsed_time,
          #  charts_generated=self._count_charts_generated(df)
     #   )
        
        print(f"✓ Presentation created from {self.data_format}: {output_file}")
       # print(f"  Processing time: {perf_summary['processing_metrics']['processing_time_seconds']}s")
       # print(f"  Rows/sec: {perf_summary['processing_metrics']['rows_processed_per_second']}")
        
        return True
    
    def _count_charts_generated(self, df):
        """Count number of charts that will be generated"""
        numeric_cols = get_numeric_columns(df)
        categorical_cols = get_categorical_columns(df)
        count = 0
        
        if numeric_cols:
            count += 1  # histogram
        if categorical_cols:
            count += 1  # bar chart
        if len(numeric_cols) >= 2:
            count += 1  # scatter plot
        
        return count
    
    def _create_presentation(self, df):
        """Create complete PowerPoint presentation"""
        self.prs = Presentation()
        self.prs.slide_width = Inches(10)
        self.prs.slide_height = Inches(7.5)
        
        self._add_title_slide()
        self._add_overview_slide(df)
        self._add_data_preview_slide(df)
        self._add_data_summary_slide(df)
        self._add_statistics_slides(df)
        from data_canvas.ai_insights import generate_insights
        summary_text = str(df.describe())
        insights = generate_insights(summary_text)
        chunks = self._split_text(insights)

        for chunk in chunks:
            self._add_ai_insights_slide(chunk)
      #  self._add_visualization_slides(df)
    
    def _add_title_slide(self):
        """Add professional title slide"""
        slide_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(slide_layout)
        
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(31, 78, 121)
        
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(9), Inches(2))
        title_frame = title_box.text_frame
        title_frame.word_wrap = True
        title_para = title_frame.paragraphs[0]
        title_para.text = self.title
        title_para.font.size = Pt(54)
        title_para.font.bold = True
        title_para.font.color.rgb = RGBColor(255, 255, 255)
        title_para.alignment = PP_ALIGN.CENTER
        
        if self.data_format:
            subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(4), Inches(9), Inches(0.8))
            subtitle_frame = subtitle_box.text_frame
            subtitle_para = subtitle_frame.paragraphs[0]
            subtitle_para.text = f"Data Format: {self.data_format}"
            subtitle_para.font.size = Pt(24)
            subtitle_para.font.color.rgb = RGBColor(200, 200, 200)
            subtitle_para.alignment = PP_ALIGN.CENTER
        
        date_box = slide.shapes.add_textbox(Inches(0.5), Inches(5.5), Inches(9), Inches(1))
        date_frame = date_box.text_frame
        date_para = date_frame.paragraphs[0]
        date_para.text = f"Generated on {datetime.now().strftime('%B %d, %Y at %H:%M')}"
        date_para.font.size = Pt(18)
        date_para.font.color.rgb = RGBColor(255, 255, 255)
        date_para.alignment = PP_ALIGN.CENTER
    
    def _add_overview_slide(self, df):
        """Add data overview slide"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[5])
        
        title = slide.shapes.title
        title.text = "📊 Data Overview"
        title.text_frame.paragraphs[0].font.size = Pt(40)
        title.text_frame.paragraphs[0].font.bold = True
        
        analyzer = DataAnalyzer(df)
        overview = analyzer.get_overview()
        
        content_box = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(5.5))
        text_frame = content_box.text_frame
        text_frame.word_wrap = True
        
        stats = [
            (f"📈 Total Rows", f"{overview['total_rows']:,}"),
            (f"📋 Total Columns", f"{overview['total_columns']}"),
            (f"🔢 Numeric Columns", f"{overview['numeric_columns']}"),
            (f"🏷️  Categorical Columns", f"{overview['categorical_columns']}"),
            (f"📁 Data Source Format", f"{self.data_format or 'Unknown'}"),
        ]
        
        for i, (label, value) in enumerate(stats):
            if i == 0:
                p = text_frame.paragraphs[0]
            else:
                p = text_frame.add_paragraph()
            
            p.text = f"{label}: {value}"
            p.font.size = Pt(22)
            p.space_before = Pt(16)
            p.level = 0
    
    def _add_data_preview_slide(self, df):
        """Add first few rows preview slide"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[5])
        
        title = slide.shapes.title
        title.text = "👁️ Data Preview (First 5 Rows)"
        title.text_frame.paragraphs[0].font.size = Pt(40)
        
        preview_df = df.head(5)
        rows, cols = min(len(preview_df) + 1, 6), min(len(preview_df.columns), 5)
        
        left = Inches(0.5)
        top = Inches(1.5)
        width = Inches(9)
        height = Inches(5.5)
        
        table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
        table = table_shape.table
        
        for col in range(cols):
            table.columns[col].width = Inches(9 / cols)
        
        for col_idx in range(cols):
            cell = table.cell(0, col_idx)
            cell.text = str(preview_df.columns[col_idx])
            cell.text_frame.paragraphs[0].font.bold = True
            cell.text_frame.paragraphs[0].font.size = Pt(11)
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(31, 78, 121)
            cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        
        for row_idx in range(len(preview_df)):
            for col_idx in range(cols):
                cell = table.cell(row_idx + 1, col_idx)
                cell.text = str(preview_df.iloc[row_idx, col_idx])
                cell.text_frame.paragraphs[0].font.size = Pt(10)
    
    def _add_data_summary_slide(self, df):
        """Add detailed data summary slide"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[5])
        
        title = slide.shapes.title
        title.text = "📋 Column Summary"
        title.text_frame.paragraphs[0].font.size = Pt(40)
        
        rows, cols = min(len(df.columns) + 1, 10), 4
        left = Inches(0.5)
        top = Inches(1.5)
        width = Inches(9)
        height = Inches(5.5)
        
        table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
        table = table_shape.table
        
        for col in range(cols):
            table.columns[col].width = Inches(9 / cols)
        
        headers = ["Column Name", "Data Type", "Non-Null", "Unique Values"]
        for col_idx, header in enumerate(headers):
            cell = table.cell(0, col_idx)
            cell.text = header
            cell.text_frame.paragraphs[0].font.bold = True
            cell.text_frame.paragraphs[0].font.size = Pt(11)
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(31, 78, 121)
            cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        
        for row_idx, col_name in enumerate(df.columns[:rows-1], 1):
            table.cell(row_idx, 0).text = str(col_name)
            table.cell(row_idx, 1).text = str(df[col_name].dtype)
            table.cell(row_idx, 2).text = str(df[col_name].notna().sum())
            table.cell(row_idx, 3).text = str(df[col_name].nunique())
            
            for col_idx in range(cols):
                table.cell(row_idx, col_idx).text_frame.paragraphs[0].font.size = Pt(10)
    
    def _add_statistics_slides(self, df):
        """Add statistics slides for numeric columns"""
        analyzer = DataAnalyzer(df)
        numeric_cols = get_numeric_columns(df)
        
        for col in numeric_cols[:5]:
            slide = self.prs.slides.add_slide(self.prs.slide_layouts[5])
            
            title = slide.shapes.title
            title.text = f"📈 Statistics: {col}"
            title.text_frame.paragraphs[0].font.size = Pt(36)
            
            stats = analyzer.get_descriptive_stats(col)
            
            content_box = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(5.5))
            text_frame = content_box.text_frame
            text_frame.word_wrap = True
            
            stats_display = [
                (f"Mean", f"{stats['mean']:.2f}"),
                (f"Median", f"{stats['median']:.2f}"),
                (f"Std Dev", f"{stats['std_dev']:.2f}"),
                (f"Min Value", f"{stats['min']:.2f}"),
                (f"Max Value", f"{stats['max']:.2f}"),
                #(f"Q1 (25%)", f"{stats.get['q25', 0]:.2f}"),
              #  (f"Q3 (75%)", f"{stats.get['q75', 0]:.2f}"),
            ]
            
            for i, (label, value) in enumerate(stats_display):
                if i == 0:
                    p = text_frame.paragraphs[0]
                else:
                    p = text_frame.add_paragraph()
                
                p.text = f"{label}: {value}"
                p.font.size = Pt(20)
                p.space_before = Pt(10)
                p.level = 0

    def _split_text(self, text, size =500):
        return [text[i:i+size] for i in range (0, len (text),size)]
        
    def _add_ai_insights_slide(self, insights):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[5])

     # Background color
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(230, 240, 255)

        title = slide.shapes.title
        title.text = "AI Insights"

    # Title color
        title.text_frame.paragraphs[0].font.color.rgb = RGBColor(31, 78, 121)
    
        
        
        textbox = slide.shapes.add_textbox( Inches (0.1), Inches (1.2), Inches (9.6), Inches (5.8)) 

        # Textbox color
        textbox.fill.solid()
        textbox.fill.fore_color.rgb = RGBColor(245, 250, 255)

        text_frame = textbox.text_frame
        text_frame.word_wrap = True 
        insights = insights.replace("###","")
        insights = insights.replace("**","")
        text_frame.text = insights
        for paragraph in text_frame.paragraphs: 
            paragraph.font.size = Pt(8)
    
    def _add_visualization_slides(self, df):
        """Add chart visualization slides"""
        numeric_cols = get_numeric_columns(df)
        categorical_cols = get_categorical_columns(df)
        
        # Downsample if needed for visualization
      #  if self.enable_optimization:
       #     viz_df = self.performance_optimizer.downsample_for_visualization(df)
     #   else:
      #      viz_df = df
        
        if numeric_cols:
            slide = self.prs.slides.add_slide(self.prs.slide_layouts[5])
            title = slide.shapes.title
            title.text = f"📊 Distribution: {numeric_cols[0]}"
            title.text_frame.paragraphs[0].font.size = Pt(36)
            
           # chart_img = self.chart_generator.histogram(
              #  viz_df[numeric_cols[0]],
               # title="Histogram",
              #  bins=20
         #   )
            
         #   self._add_image_to_slide(slide, chart_img)
        
        if categorical_cols:
            slide = self.prs.slides.add_slide(self.prs.slide_layouts[5])
            title = slide.shapes.title
            title.text = f"🏷️  Top Values: {categorical_cols[0]}"
            title.text_frame.paragraphs[0].font.size = Pt(36)
            
         #   top_vals = viz_df[categorical_cols[0]].value_counts().head(10)
         #   chart_img = self.chart_generator.bar_chart(
          #      top_vals,
           #     title="Bar Chart"
        #    )
            
          #  self._add_image_to_slide(slide, chart_img)
        
        if len(numeric_cols) >= 2:
            slide = self.prs.slides.add_slide(self.prs.slide_layouts[5])
            title = slide.shapes.title
            title.text = f"🔗 Relationship: {numeric_cols[0]} vs {numeric_cols[1]}"
            title.text_frame.paragraphs[0].font.size = Pt(36)
            
         #   chart_img = self.chart_generator.scatter_plot(
            #    viz_df,
            #    numeric_cols[0],
            #    numeric_cols[1],
             #   title="Scatter Plot"
         #   )
           
         #   self._add_image_to_slide(slide, chart_img)
    
    def _add_image_to_slide(self, slide, image, left=Inches(1), top=Inches(1.5)):
        """Add image to slide"""
        img_bytes = io.BytesIO()
        image.save(img_bytes, format="PNG")
        img_bytes.seek(0)
        
        slide.shapes.add_picture(img_bytes, left, top, width=Inches(8))
    
    #def get_performance_metrics(self):
        #"""Get performance metrics from conversion"""
        #return self.performance_optimizer.get_performance_summary()
    
    def get_structured_report(self):
        """Get structured report from conversion"""
        return self.structured_report
